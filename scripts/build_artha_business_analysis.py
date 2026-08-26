from __future__ import annotations

import math
import os
from pathlib import Path

from docx import Document
from docx.enum.section import WD_ORIENT
from docx.enum.table import WD_ALIGN_VERTICAL, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PPTColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path("/Users/sintzu/TaxDisputeC")
OUT = ROOT / "outputs" / "artha-ai-jurist-analysis"
OUT.mkdir(parents=True, exist_ok=True)

DOCX_OUT = OUT / "Artha_AI_Jurist_Business_Analysis_Report_ID.docx"
PPTX_OUT = OUT / "Artha_AI_Jurist_Refined_Leadership_Deck.pptx"

RSM_BLUE = "00A3E0"
RSM_GREEN = "3FA044"
RSM_GRAY = "555A5D"
RSM_DARK = "2B2F3A"
RSM_LIGHT_BLUE = "EAF7FC"
RSM_LIGHT_GREEN = "EAF6EC"
RSM_BORDER = "D8E0E6"
RSM_MUTED = "667085"


YEARS = ["FY1", "FY2", "FY3", "FY4", "FY5"]
REVENUE = [564, 2160, 6840, 15840, 30000]  # Rp mm
COGS = [250, 800, 1950, 3800, 6000]
GP = [r - c for r, c in zip(REVENUE, COGS)]
GM = [g / r for g, r in zip(GP, REVENUE)]
OPEX = [2400, 3010, 3942, 4147.2, 6482.48]
EBITDA = [g - o for g, o in zip(GP, OPEX)]
INITIAL_CAPITAL = 4000

PRICING = [
    ("Lite", "Students / research funnel", 0.6, "low-price acquisition layer"),
    ("Premium", "Tax professionals / corporate tax users", 4.8, "core paid research tier"),
    ("Workflow", "Advisors handling disputes", 24.0, "drafting and case-preparation layer"),
    ("Enterprise", "Corporate account / RSM license", 120.0, "account-level deployment"),
]

VOLUMES = {
    "Lite": [500, 1000, 2000, 4000, 8000],
    "Premium": [50, 150, 500, 1000, 2000],
    "Workflow": [1, 25, 100, 280, 500],
    "Enterprise": [0, 2, 7, 16, 30],
}

UNIT_ECON = [
    ("Lite", 0.6, 0.3, "35%", "70%", 1.2, "4.0x", "8.6 months"),
    ("Premium", 4.8, 3.0, "20%", "75%", 18.0, "6.0x", "10.0 months"),
    ("Workflow", 24.0, 15.0, "12%", "78%", 156.0, "10.4x", "9.6 months"),
    ("Enterprise", 120.0, 50.0, "8%", "78%", 1170.0, "23.4x", "6.4 months"),
]


def npv(rate: float, cashflows: list[float]) -> float:
    return sum(cf / ((1 + rate) ** i) for i, cf in enumerate(cashflows))


def irr(cashflows: list[float]) -> float:
    lo, hi = -0.99, 5.0
    for _ in range(200):
        mid = (lo + hi) / 2
        if npv(mid, cashflows) > 0:
            lo = mid
        else:
            hi = mid
    return (lo + hi) / 2


def payback(cashflows: list[float]) -> float | None:
    cumulative = cashflows[0]
    for i in range(1, len(cashflows)):
        previous = cumulative
        cumulative += cashflows[i]
        if cumulative >= 0:
            return (i - 1) + (-previous / cashflows[i])
    return None


def rp_bn(mm: float) -> str:
    return f"Rp{mm / 1000:,.1f}bn".replace(",", ".")


def rp_mm(mm: float) -> str:
    if abs(mm) >= 1000:
        return rp_bn(mm)
    if abs(mm) < 10 and abs(mm - round(mm)) > 0.001:
        return f"Rp{mm:.1f}m"
    return f"Rp{mm:,.0f}m".replace(",", ".")


def pct(x: float, digits: int = 1) -> str:
    return f"{x * 100:.{digits}f}%"


BASE_CASHFLOWS = [-INITIAL_CAPITAL] + EBITDA
BEP_REVENUE = [o / gm for o, gm in zip(OPEX, GM)]
SCENARIO_MULTIPLIERS = {
    "Worst-case": 0.6,
    "Basic": 0.8,
    "Baseline": 1.0,
    "Moderate": 1.4,
}


def scenario_rows():
    rows = []
    for name, multiplier in SCENARIO_MULTIPLIERS.items():
        scenario_revenue = [r * multiplier for r in REVENUE]
        scenario_gp = [scenario_revenue[i] * GM[i] for i in range(5)]
        scenario_ebitda = [scenario_gp[i] - OPEX[i] for i in range(5)]
        cashflows = [-INITIAL_CAPITAL] + scenario_ebitda
        rows.append(
            {
                "name": name,
                "multiplier": multiplier,
                "fy5_revenue": scenario_revenue[-1],
                "fy5_ebitda": scenario_ebitda[-1],
                "npv15": npv(0.15, cashflows),
                "irr": irr(cashflows),
                "payback": payback(cashflows),
            }
        )
    return rows


def set_cell_shading(cell, fill: str):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(qn("w:shd"))
    if shd is None:
        shd = OxmlElement("w:shd")
        tc_pr.append(shd)
    shd.set(qn("w:fill"), fill)


def set_cell_text(cell, text: str, bold: bool = False, color: str = RSM_DARK, size: int = 9):
    cell.text = ""
    p = cell.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run(str(text))
    run.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor.from_string(color)
    cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER


def set_table_widths(table, widths):
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for row in table.rows:
        for idx, width in enumerate(widths):
            if idx < len(row.cells):
                row.cells[idx].width = Inches(width)


def add_table(doc, headers, rows, widths=None, font_size=8):
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    for i, header in enumerate(headers):
        set_cell_text(hdr[i], header, bold=True, color="FFFFFF", size=font_size)
        set_cell_shading(hdr[i], RSM_GRAY)
    for row in rows:
        cells = table.add_row().cells
        for i, value in enumerate(row):
            set_cell_text(cells[i], value, size=font_size)
            if len(table.rows) % 2 == 0:
                set_cell_shading(cells[i], "F7FAFC")
    if widths:
        set_table_widths(table, widths)
    doc.add_paragraph()
    return table


def add_callout(doc, title: str, body: str, fill: str = RSM_LIGHT_GREEN):
    table = doc.add_table(rows=1, cols=1)
    table.style = "Table Grid"
    cell = table.cell(0, 0)
    set_cell_shading(cell, fill)
    p = cell.paragraphs[0]
    p.paragraph_format.space_after = Pt(4)
    r = p.add_run(title)
    r.bold = True
    r.font.color.rgb = RGBColor.from_string(RSM_DARK)
    r.font.size = Pt(10)
    p2 = cell.add_paragraph()
    p2.paragraph_format.space_after = Pt(0)
    r2 = p2.add_run(body)
    r2.font.size = Pt(9)
    r2.font.color.rgb = RGBColor.from_string(RSM_DARK)
    doc.add_paragraph()


def add_bullets(doc, items):
    for item in items:
        p = doc.add_paragraph(style="List Bullet")
        run = p.add_run(item)
        run.font.size = Pt(10)


def setup_doc_styles(doc: Document):
    section = doc.sections[0]
    section.orientation = WD_ORIENT.PORTRAIT
    section.top_margin = Inches(0.8)
    section.bottom_margin = Inches(0.75)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)

    styles = doc.styles
    normal = styles["Normal"]
    normal.font.name = "Aptos"
    normal.font.size = Pt(10)
    normal.font.color.rgb = RGBColor.from_string(RSM_DARK)
    normal.paragraph_format.space_after = Pt(6)
    normal.paragraph_format.line_spacing = 1.1

    for name, size, color in [
        ("Heading 1", 16, RSM_BLUE),
        ("Heading 2", 13, RSM_BLUE),
        ("Heading 3", 11, RSM_DARK),
    ]:
        style = styles[name]
        style.font.name = "Aptos Display"
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = RGBColor.from_string(color)
        style.paragraph_format.space_before = Pt(12)
        style.paragraph_format.space_after = Pt(6)


def build_docx():
    doc = Document()
    setup_doc_styles(doc)

    # Cover
    p = doc.add_paragraph()
    r = p.add_run("RSM")
    r.bold = True
    r.font.size = Pt(36)
    r.font.color.rgb = RGBColor.from_string(RSM_GRAY)
    p.paragraph_format.space_after = Pt(2)
    p2 = doc.add_paragraph()
    r2 = p2.add_run("Artha AI Jurist / TaxCo Business Analysis")
    r2.bold = True
    r2.font.size = Pt(22)
    r2.font.color.rgb = RGBColor.from_string(RSM_DARK)
    p3 = doc.add_paragraph()
    p3.add_run("Senior business analysis, financial validation, external benchmarks, and leadership-deck recommendations").font.size = Pt(11)
    p3.runs[0].font.color.rgb = RGBColor.from_string(RSM_MUTED)
    add_callout(
        doc,
        "Executive answer",
        "The concept is commercially promising if positioned as an agentic tax-dispute workflow product rather than a generic tax library. The baseline model produces FY5 revenue of Rp30.0bn, FY5 EBITDA of Rp17.5bn, NPV at 15% of Rp6.8bn, IRR of 36.1%, and payback of approximately 3.86 years. However, the current pack needs three fixes before leadership approval: repair broken scenario formulas, reconcile the capital ask, and explicitly govern RSM data/IP, privacy, and AI quality risk.",
    )

    doc.add_heading("1. Pemahaman Dokumen", level=1)
    doc.add_paragraph(
        "Dokumen menggambarkan rencana membangun platform AI pajak Indonesia bernama Artha AI Jurist / TaxCo, dimulai dari pemanfaatan internal RSM lalu berkembang ke riset berbayar, workflow sengketa, dan enterprise licensing. Ruang lingkupnya mencakup database putusan, database peraturan, AI drafting untuk dispute, client alerts, dan treaty/MLI synthesizer. Asumsi utama: RSM menjadi anchor customer/data flywheel, monetisasi dilakukan melalui SaaS tahunan dan enterprise account, serta pasar awal berasal dari konsultan pajak, tax professional korporasi, dan segmen pendidikan."
    )

    add_table(
        doc,
        ["Kategori", "Data yang tersedia", "Implikasi bisnis"],
        [
            ["Timeline", "Internal pilot 0-6 bulan; research launch 6-12 bulan; workflow launch 12-24 bulan", "Roadmap logis, tetapi butuh gate KPI per fase."],
            ["Target pasar", "7.500 konsultan pajak; 33.300 in-house corporate seats; 40.800 professional paying universe; 50.000 students/academics", "TAM cukup besar, tetapi angka harus divalidasi dengan sumber primer."],
            ["Pricing", "Lite Rp0,6m; Premium Rp4,8m; Workflow Rp24m; Enterprise Rp120m per tahun", "Pricing bertingkat cocok untuk funnel, tetapi profit mainly driven oleh workflow dan enterprise."],
            ["Revenue", "FY1 Rp564m; FY2 Rp2,16bn; FY3 Rp6,84bn; FY4 Rp15,84bn; FY5 Rp30,0bn", "Model menuntut akselerasi penjualan tinggi setelah tahun kedua."],
            ["Profitability", "EBITDA FY3 positif Rp948m; cumulative EBITDA positif FY4", "Profitability claim benar untuk EBITDA tahunan FY3, bukan cash payback total."],
            ["Unit economics", "LTV/CAC 4,0x sampai 23,4x; CAC payback 6,4-10 bulan", "Menarik jika churn/CAC dapat dibuktikan melalui pilot."],
            ["Modal", "Workbook: initial capital Rp4,0bn; deck: Rp12,3bn funding through FY3", "Gap material; perlu rekonsiliasi sebelum fundraising/internal approval."],
            ["Model issue", "Sheet Scenarios berisi formula #REF!", "Skenario harus dibangun ulang agar deck kredibel."],
        ],
        widths=[1.3, 2.6, 2.9],
    )

    add_table(
        doc,
        ["Informasi hilang / ambigu", "Asumsi analisis", "Dampak jika asumsi salah"],
        [
            ["Sumber primer angka 16k+ tax court cases dan 20% growth", "Dipakai sebagai klaim deck, tetapi ditandai perlu verifikasi", "Market urgency bisa overstated."],
            ["Rincian Rp12,3bn capital ask", "Analisis finansial memakai Rp4,0bn initial capital dari workbook", "NPV/IRR berubah jika investasi aktual lebih besar."],
            ["Kontrak anchor RSM dan revenue commitment", "Diasumsikan belum legally committed", "Revenue FY1-FY2 berisiko lebih rendah."],
            ["Legal basis penggunaan putusan, dokumen WP, dan data internal RSM", "Diasumsikan butuh data wall, anonymization, dan policy approval", "Risiko reputasi/privacy tinggi."],
            ["Churn dan CAC aktual", "Menggunakan angka workbook", "Unit economics perlu dikalibrasi dengan pilot."],
        ],
        widths=[2.4, 2.3, 2.1],
    )

    doc.add_heading("2. Analisis Bisnis", level=1)
    doc.add_heading("Framework yang dipilih", level=2)
    add_bullets(
        doc,
        [
            "SWOT digunakan untuk melihat kelayakan strategis karena produk bergantung pada trust, data proprietary, dan timing pasar AI.",
            "Business Model Canvas dipakai karena monetisasi produk mencampur SaaS, enterprise license, dan service-enabled workflow.",
            "Porter’s Five Forces relevan untuk memetakan tekanan dari tax research provider, legal/regulatory database, dan AI horizontal tools.",
            "Stakeholder dan pain point mapping penting karena adopsi bergantung pada advisor, corporate tax team, reviewer, dan governance owner.",
            "Risk assessment likelihood x impact diperlukan karena produk menyentuh dokumen sengketa, data WP, regulasi, dan output AI yang high stakes.",
        ],
    )

    add_table(
        doc,
        ["SWOT", "Poin utama", "Implikasi"],
        [
            ["Strength", "RSM memiliki credibility, domain expertise, dan akses knowledge sengketa yang sulit ditiru.", "Moat terbaik bukan library, tetapi outcome/dispute workflow."],
            ["Weakness", "Model masih early, data source belum sepenuhnya verified, dan spreadsheet skenario rusak.", "Perlu governance dan model-control sebelum komersialisasi luas."],
            ["Opportunity", "GenAI memperbesar peluang otomasi knowledge work dan tax administration bergerak digital.", "Timing mendukung productization tax advisory."],
            ["Threat", "Provider database pajak, legaltech, dan AI generic tools dapat menyerang fitur research.", "Fokus pada workflow proprietary, bukan search generik."],
        ],
        widths=[1.1, 3.0, 2.7],
    )

    add_table(
        doc,
        ["Business Model Canvas", "Desain yang disarankan"],
        [
            ["Customer segments", "Internal RSM teams, tax consultants, corporate in-house tax, students/researchers, enterprise tax departments."],
            ["Value proposition", "Reduce dispute research and drafting cycle time; improve consistency; surface comparable decisions; manage regulation context."],
            ["Channels", "RSM internal pilot, webinars/training, professional community, direct enterprise sales, partnership with tax education providers."],
            ["Revenue streams", "Annual seats, enterprise accounts, workflow module, implementation/custom data projects, internal license."],
            ["Key resources", "Curated decision database, regulation knowledge base, extraction QA, RSM tax experts, secure cloud/data governance."],
            ["Key activities", "Data ingestion, labeling, RAG retrieval tuning, drafting QA, customer success, regulatory update curation."],
            ["Cost structure", "Engineering, AI/API compute, data processing, expert review, cloud/security, sales/marketing."],
        ],
        widths=[1.8, 5.0],
    )

    add_table(
        doc,
        ["Force", "Assessment", "Management action"],
        [
            ["Rivalry", "Medium-high: tax/legal databases already exist, but workflow AI remains nascent.", "Differentiate through dispute drafting and RSM-reviewed comparators."],
            ["Threat of substitutes", "High: users can combine ChatGPT, public decisions, and manual research.", "Offer verified sources, repeatable scoring, and exportable workpapers."],
            ["Buyer power", "Medium: professionals compare with DDTC/Hukumonline/Ortax pricing.", "Use freemium/Lite as funnel and prove time-saving ROI for premium tiers."],
            ["Supplier power", "Medium: dependent on LLM providers, cloud, and source data availability.", "Keep model abstraction layer and local structured database."],
            ["New entrants", "Medium: app development is easier, but curated case data and tax expertise are hard.", "Build proprietary labeled outcome dataset and reviewer workflow."],
        ],
        widths=[1.4, 3.0, 2.4],
    )

    add_table(
        doc,
        ["Stakeholder", "Pain point", "What the product must prove"],
        [
            ["RSM tax advisor", "Manual search, inconsistent draft quality, repetitive report preparation.", "Shorter research time and better first draft quality."],
            ["Partner / reviewer", "Quality risk and liability from AI output.", "Traceable sources, confidence scoring, and review controls."],
            ["Corporate tax team", "Need practical answer under audit/dispute deadlines.", "Clear recommendations, evidence checklist, and comparable decisions."],
            ["IT / governance", "Data privacy, access control, and vendor risk.", "Data wall, audit log, role-based access, secure storage."],
            ["Finance / leadership", "Uncertain ROI and funding ask.", "Pilot KPIs, scenario model, payback, and staged funding gates."],
        ],
        widths=[1.5, 2.7, 2.6],
    )

    add_table(
        doc,
        ["Risk", "Likelihood", "Impact", "Score", "Mitigation"],
        [
            ["Data privacy / taxpayer confidentiality", "High", "High", "9", "Anonymization, permissioning, audit log, private datasets, PDP compliance review."],
            ["AI hallucination or wrong legal inference", "Medium", "High", "6", "RAG-only citations, mandatory human review, confidence thresholds, answer disclaimers."],
            ["Source-data quality gaps", "High", "Medium", "6", "Extraction QA, reviewer validation, provenance tracking, re-extraction workflow."],
            ["Weak paid conversion from Lite", "Medium", "Medium", "4", "Gate Lite features; drive Premium/Workflow conversion through templates and draft exports."],
            ["Capital ask inconsistency", "Medium", "High", "6", "Rebuild funding plan with explicit hiring, compute, security, and buffer assumptions."],
            ["Competitive response by incumbents", "Medium", "Medium", "4", "Move first on dispute workflow and proprietary outcome database."],
        ],
        widths=[2.0, 0.8, 0.8, 0.6, 2.6],
    )

    doc.add_heading("3. Perhitungan Kuantitatif", level=1)
    add_callout(
        doc,
        "Basis perhitungan",
        "Semua angka finansial memakai satuan Rp juta, kecuali disebut lain. Discount rate baseline 15% digunakan sebagai proxy hurdle rate early-stage SaaS/AI venture. Initial capital memakai Rp4,0bn dari workbook, bukan Rp12,3bn di deck lama karena rinciannya belum tersedia.",
        fill=RSM_LIGHT_BLUE,
    )

    add_table(
        doc,
        ["Formula", "Rumus eksplisit"],
        [
            ["Revenue", "Units x annual price per tier; scenario revenue = baseline revenue x scenario multiplier."],
            ["Gross profit", "Revenue - COGS."],
            ["Gross margin", "Gross profit / Revenue."],
            ["EBITDA", "Gross profit - Total opex."],
            ["NPV", "-Initial capital + sum(EBITDA_t / (1 + discount rate)^t)."],
            ["IRR", "Discount rate that makes NPV = 0."],
            ["Payback", "Year before cumulative cash flow turns positive + remaining deficit / next-year cash flow."],
            ["ROI", "Sum of 5Y EBITDA / initial capital."],
            ["Break-even revenue", "Fixed opex / gross margin."],
        ],
        widths=[1.7, 5.1],
    )

    add_table(
        doc,
        ["Metric"] + YEARS,
        [
            ["Revenue"] + [rp_mm(x) for x in REVENUE],
            ["COGS"] + [rp_mm(x) for x in COGS],
            ["Gross profit"] + [rp_mm(x) for x in GP],
            ["Gross margin"] + [pct(x) for x in GM],
            ["Opex"] + [rp_mm(x) for x in OPEX],
            ["EBITDA"] + [rp_mm(x) for x in EBITDA],
            ["EBITDA margin"] + [pct(EBITDA[i] / REVENUE[i]) for i in range(5)],
            ["Cumulative EBITDA"] + [rp_mm(sum(EBITDA[: i + 1])) for i in range(5)],
            ["Break-even revenue"] + [rp_mm(x) for x in BEP_REVENUE],
        ],
        widths=[1.4, 1.05, 1.05, 1.05, 1.05, 1.05],
        font_size=7,
    )

    add_table(
        doc,
        ["Financial indicator", "Result", "Interpretation"],
        [
            ["NPV @ 12%", rp_bn(npv(0.12, BASE_CASHFLOWS)), "Positive with meaningful buffer."],
            ["NPV @ 15%", rp_bn(npv(0.15, BASE_CASHFLOWS)), "Baseline value creation remains positive."],
            ["NPV @ 20%", rp_bn(npv(0.20, BASE_CASHFLOWS)), "Still positive, but more sensitive to execution."],
            ["IRR", pct(irr(BASE_CASHFLOWS), 1), "Attractive for a high-risk early venture if assumptions hold."],
            ["Payback", f"{payback(BASE_CASHFLOWS):.2f} years", "Capital recovered during FY4 under baseline assumptions."],
            ["5Y EBITDA ROI", f"{sum(EBITDA) / INITIAL_CAPITAL:.2f}x", "Strong, but depends on FY4-FY5 scale-up."],
        ],
        widths=[1.8, 1.4, 3.6],
    )

    add_table(
        doc,
        ["Scenario", "Volume/revenue multiplier", "FY5 revenue", "FY5 EBITDA", "NPV @15%", "IRR", "Payback"],
        [
            [
                r["name"],
                f"{r['multiplier']:.1f}x",
                rp_mm(r["fy5_revenue"]),
                rp_mm(r["fy5_ebitda"]),
                rp_mm(r["npv15"]),
                pct(r["irr"], 1),
                f"{r['payback']:.2f} yrs" if r["payback"] else "n/a",
            ]
            for r in scenario_rows()
        ],
        widths=[1.2, 1.2, 1.1, 1.1, 1.1, 0.8, 0.9],
        font_size=7,
    )

    add_table(
        doc,
        ["Tier", "ARPU/yr", "CAC", "Churn", "GM", "LTV", "LTV/CAC", "CAC payback"],
        [[t, rp_mm(a), rp_mm(c), ch, gm, rp_mm(ltv), ratio, pay] for t, a, c, ch, gm, ltv, ratio, pay in UNIT_ECON],
        widths=[1.1, 0.8, 0.8, 0.8, 0.7, 0.9, 0.8, 1.0],
        font_size=7,
    )

    doc.add_heading("4. Riset Referensi Eksternal", level=1)
    add_table(
        doc,
        ["Source", "Tanggal / tipe", "Klaim yang dipakai", "Status"],
        [
            ["OECD, Tax Administration 3.0: The Digital Transformation of Tax Administration", "2020, laporan resmi", "Tax administration is moving toward more integrated, digital, event-based systems; supports urgency for digital tax workflow.", "Fakta eksternal"],
            ["McKinsey, The economic potential of generative AI: The next productivity frontier", "14 Jun 2023, industry report", "GenAI has large productivity potential in knowledge work; validates opportunity for AI-assisted research/drafting.", "Benchmark eksternal"],
            ["NBER Working Paper 31161, Generative AI at Work", "Apr 2023, working paper", "AI assistance improved productivity in a real customer-support setting, with larger benefit for less-experienced workers.", "Evidence eksternal"],
            ["UU No. 27 Tahun 2022 tentang Perlindungan Data Pribadi, BPK/JDIH", "2022, regulasi", "Personal data handling creates compliance obligations; relevant to taxpayer documents and user logs.", "Regulatory fact"],
            ["TaxCo_Model.xlsx", "Dokumen internal", "TAM, pricing benchmark, volume forecast, and unit economics.", "Model estimate; perlu validasi primer"],
        ],
        widths=[2.2, 1.3, 2.4, 0.9],
        font_size=7,
    )
    doc.add_paragraph("Reference URLs:")
    add_bullets(
        doc,
        [
            "OECD Tax Administration 3.0: https://www.oecd.org/tax/forum-on-tax-administration/publications-and-products/tax-administration-3-0-the-digital-transformation-of-tax-administration.htm",
            "McKinsey GenAI productivity report: https://www.mckinsey.com/capabilities/mckinsey-digital/our-insights/the-economic-potential-of-generative-ai-the-next-productivity-frontier",
            "NBER Generative AI at Work: https://www.nber.org/papers/w31161",
            "UU PDP No. 27/2022: https://peraturan.bpk.go.id/Details/229798/uu-no-27-tahun-2022",
        ],
    )

    doc.add_heading("5. Temuan dan Rekomendasi", level=1)
    add_table(
        doc,
        ["Priority", "Finding", "Business impact", "Recommendation"],
        [
            ["High", "Positioning as library alone is weak.", "Competes with DDTC/Hukumonline/Ortax and generic AI.", "Lead with Dispute Co-pilot + verified comparator database."],
            ["High", "Financial story attractive but deck has inconsistencies.", "Leadership may question credibility.", "Fix scenario formulas, reconcile Rp4bn vs Rp12.3bn, and clarify EBITDA vs gross margin."],
            ["High", "RSM anchor/data flywheel is the real moat.", "Creates defensibility and quality loop.", "Make internal pilot mandatory, with reviewer feedback feeding the model."],
            ["Medium", "Lite is useful as funnel but not profit engine.", "Volume-heavy, low ARPU, support burden risk.", "Use Lite for lead generation; push Premium/Workflow conversion."],
            ["Medium", "Governance/security is not optional.", "Taxpayer data and AI advice are high-stakes.", "Build role-based access, audit log, data masking, source traceability, and human approval gates."],
            ["Low", "Market claim needs source hardening.", "Weakens investor/leadership confidence.", "Add verified sources for tax court case volume, consultant count, and corporate SPT universe."],
        ],
        widths=[0.8, 1.8, 1.8, 2.4],
        font_size=7,
    )

    add_table(
        doc,
        ["Quick wins (0-90 days)", "Longer-term initiatives"],
        [
            ["Repair spreadsheet scenario sheet and add an assumption register.", "Build proprietary labeled dispute-outcome database and reviewer tooling."],
            ["Run RSM internal pilot with measured time-saving and draft-quality KPIs.", "Launch enterprise-grade security, data wall, and model evaluation framework."],
            ["Rewrite deck slides 9-10 with corrected EBITDA, NPV, IRR, payback, and capital ask.", "Create API/data architecture for regulation updates and firm knowledge ingestion."],
            ["Validate WTP with 20-30 advisor/corporate interviews.", "Develop partner channels with tax education and professional associations."],
        ],
        widths=[3.3, 3.5],
    )

    doc.add_heading("Slide Improvement Notes", level=1)
    add_bullets(
        doc,
        [
            "Change the main story from 'AI tax library' to 'agentic tax-dispute workflow and drafting platform'.",
            "Add a one-slide financial scorecard with NPV, IRR, payback, FY5 revenue, FY5 EBITDA, and break-even year.",
            "Replace the broken scenario sheet with four transparent cases: Worst-case, Basic, Baseline, Moderate.",
            "Add an explicit decision ask: approve 90-day pilot, approve model rebuild, and decide capital envelope after KPI review.",
            "Move market-size claims to sourced footnotes and mark assumptions that still require validation.",
            "Add governance slide covering data privacy, RAG citations, human review, audit logs, and data wall.",
        ],
    )

    doc.add_heading("Pertanyaan Klarifikasi", level=1)
    add_bullets(
        doc,
        [
            "Apakah angka kebutuhan modal yang benar Rp4,0bn, Rp12,3bn, atau kombinasi initial capital + buffer + hiring plan?",
            "Apakah RSM akan menjadi paying anchor customer atau hanya pilot sponsor?",
            "Apakah data sengketa internal RSM dapat digunakan untuk training/RAG setelah anonymization dan client consent review?",
            "Berapa target waktu saving yang ingin dibuktikan dalam pilot: 30%, 50%, atau lebih?",
            "Apakah produk akan dijual sebagai separate entity, RSM-branded tool, atau white-label internal platform?",
            "Sumber primer mana yang akan dipakai untuk klaim tax court case volume dan growth?",
        ],
    )

    footer = doc.sections[0].footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fr = footer.add_run("RSM / Artha AI Jurist business analysis - draft for leadership discussion")
    fr.font.size = Pt(8)
    fr.font.color.rgb = RGBColor.from_string(RSM_MUTED)
    doc.save(DOCX_OUT)


# ---------------------- PPTX ----------------------


def rgb(hex_color: str) -> PPTColor:
    return PPTColor.from_string(hex_color.replace("#", ""))


def add_text(slide, text, x, y, w, h, size=20, bold=False, color=RSM_DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = PptInches(0.02)
    tf.margin_right = PptInches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = PptPt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return box


def add_small_footer(slide, text="Source: TaxCo_Model.xlsx and Artha AI Jurist deck; external benchmarks cited in report."):
    add_text(slide, text, 0.55, 7.05, 12.2, 0.22, size=7, color=RSM_MUTED)


def add_rsm_logo(slide, x=0.55, y=0.3, scale=1.0):
    # RSM bar motif, editable and lightweight.
    h = 0.08 * scale
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(x), PptInches(y), PptInches(0.18 * scale), PptInches(h)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = rgb("8A9299")
    slide.shapes[-1].line.fill.background()
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(x + 0.28 * scale), PptInches(y), PptInches(0.55 * scale), PptInches(h)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = rgb(RSM_GREEN)
    slide.shapes[-1].line.fill.background()
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(x + 0.95 * scale), PptInches(y), PptInches(1.05 * scale), PptInches(h)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = rgb(RSM_BLUE)
    slide.shapes[-1].line.fill.background()
    add_text(slide, "RSM", x, y + 0.17 * scale, 1.45 * scale, 0.48 * scale, size=int(28 * scale), bold=True, color=RSM_GRAY)


def add_title(slide, title, subtitle=None):
    add_rsm_logo(slide)
    add_text(slide, title, 0.55, 0.88, 12.0, 0.62, size=24, bold=True, color=RSM_DARK)
    if subtitle:
        add_text(slide, subtitle, 0.58, 1.52, 11.8, 0.34, size=10, color=RSM_MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0.55), PptInches(1.9), PptInches(12.2), PptInches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(RSM_BORDER)
    line.line.fill.background()


def add_card(slide, x, y, w, h, title, body="", color=RSM_BLUE, title_size=15):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("FFFFFF")
    shape.line.color.rgb = rgb(RSM_BORDER)
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(0.05))
    top.fill.solid()
    top.fill.fore_color.rgb = rgb(color)
    top.line.fill.background()
    add_text(slide, title, x + 0.18, y + 0.18, w - 0.36, 0.35, size=title_size, bold=True, color=RSM_DARK)
    if body:
        add_text(slide, body, x + 0.18, y + 0.58, w - 0.36, h - 0.7, size=10, color=RSM_MUTED)
    return shape


def add_metric(slide, x, y, w, h, label, value, note="", color=RSM_BLUE):
    add_card(slide, x, y, w, h, label, "", color=color, title_size=10)
    add_text(slide, value, x + 0.18, y + 0.55, w - 0.3, 0.48, size=23, bold=True, color=RSM_GRAY)
    if note:
        add_text(slide, note, x + 0.18, y + 1.08, w - 0.3, 0.32, size=8, color=RSM_MUTED)


def add_bar(slide, x, y, label, value, max_value, color, suffix=""):
    add_text(slide, label, x, y - 0.03, 1.65, 0.25, size=8, bold=True, color=RSM_MUTED)
    track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(x + 1.75), PptInches(y), PptInches(3.4), PptInches(0.12))
    track.fill.solid()
    track.fill.fore_color.rgb = rgb("E6EDF2")
    track.line.fill.background()
    width = 3.4 * min(1, value / max_value)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(x + 1.75), PptInches(y), PptInches(width), PptInches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(color)
    bar.line.fill.background()
    add_text(slide, f"{value:.1f}{suffix}", x + 5.25, y - 0.06, 0.65, 0.25, size=8, bold=True, color=RSM_DARK)


def add_bullets_ppt(slide, items, x, y, w, h, size=11, color=RSM_DARK):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.margin_left = PptInches(0.15)
        p.font.name = "Aptos"
        p.font.size = PptPt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = PptPt(5)
    return box


def make_table_ppt(slide, x, y, w, h, headers, rows, col_widths=None, font_size=8):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    table = table_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = PptInches(w * cw / total)
    for r_i in range(len(rows) + 1):
        row = table.rows[r_i]
        row.height = PptInches(h / (len(rows) + 1))
        for c_i in range(len(headers)):
            cell = table.cell(r_i, c_i)
            cell.margin_left = PptInches(0.06)
            cell.margin_right = PptInches(0.06)
            cell.margin_top = PptInches(0.03)
            cell.margin_bottom = PptInches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(RSM_GRAY if r_i == 0 else ("F7FAFC" if r_i % 2 == 0 else "FFFFFF"))
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = headers[c_i] if r_i == 0 else str(rows[r_i - 1][c_i])
            r.font.name = "Aptos"
            r.font.size = PptPt(font_size)
            r.font.bold = r_i == 0
            r.font.color.rgb = rgb("FFFFFF" if r_i == 0 else RSM_DARK)
    return table_shape


def build_pptx():
    prs = Presentation()
    prs.slide_width = PptInches(13.333333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Cover
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("FFFFFF")
    add_rsm_logo(slide, 0.65, 0.55, 1.2)
    add_text(slide, "Artha AI Jurist", 0.65, 1.65, 5.2, 0.55, size=30, bold=True, color=RSM_DARK)
    add_text(slide, "Business analysis and refined leadership case", 0.67, 2.15, 5.8, 0.38, size=14, color=RSM_MUTED)
    add_card(
        slide,
        7.0,
        1.15,
        5.15,
        2.05,
        "Investment thesis",
        "The opportunity is strongest when positioned as an agentic tax-dispute workflow platform, not as another tax library.",
        color=RSM_GREEN,
        title_size=18,
    )
    add_metric(slide, 0.65, 4.05, 2.25, 1.35, "FY5 revenue", rp_bn(REVENUE[-1]), "Baseline model", RSM_BLUE)
    add_metric(slide, 3.1, 4.05, 2.25, 1.35, "FY5 EBITDA", rp_bn(EBITDA[-1]), "58.4% margin", RSM_GREEN)
    add_metric(slide, 5.55, 4.05, 2.25, 1.35, "NPV @15%", rp_bn(npv(0.15, BASE_CASHFLOWS)), "Rp4.0bn initial capital", RSM_BLUE)
    add_metric(slide, 8.0, 4.05, 2.25, 1.35, "IRR", pct(irr(BASE_CASHFLOWS), 1), "Baseline", RSM_GRAY)
    add_metric(slide, 10.45, 4.05, 2.25, 1.35, "Payback", f"{payback(BASE_CASHFLOWS):.2f} yrs", "During FY4", RSM_GREEN)
    add_small_footer(slide, "Draft for leadership discussion. Financials from TaxCo_Model.xlsx, recalculated where workbook scenarios returned #REF!.")

    # 2 Problem
    slide = prs.slides.add_slide(blank)
    add_title(slide, "1. Problem: tax dispute work is high value but still manual", "The proposed product should solve workflow pain, not only search pain.")
    add_card(slide, 0.7, 2.25, 3.75, 3.3, "Pain point", "Advisors spend time locating relevant decisions, checking regulations, aligning evidence, and drafting client-facing recommendations.", RSM_BLUE)
    add_card(slide, 4.8, 2.25, 3.75, 3.3, "Market signal", "The source deck cites 16k+ annual tax court cases and long dispute timelines. These claims should be source-hardened before external fundraising.", RSM_GREEN)
    add_card(slide, 8.9, 2.25, 3.75, 3.3, "Why AI now", "GenAI is suited to knowledge-work tasks: retrieval, summarization, issue mapping, and first-draft generation with human review.", RSM_GRAY)
    add_small_footer(slide, "External support: OECD Tax Administration 3.0; McKinsey GenAI productivity report; NBER Generative AI at Work.")

    # 3 Market
    slide = prs.slides.add_slide(blank)
    add_title(slide, "2. Addressable market: enough depth, but validate sources", "The attached model frames a professional paying universe of 40.8k plus a student/research funnel.")
    add_metric(slide, 0.7, 2.25, 2.2, 1.25, "Tax consultants", "7.5k", "Model source: SIKOP/IKPI", RSM_BLUE)
    add_metric(slide, 3.15, 2.25, 2.2, 1.25, "Corporate seats", "33.3k", "5% of 444k x 1.5 seats", RSM_GREEN)
    add_metric(slide, 5.6, 2.25, 2.2, 1.25, "Professional universe", "40.8k", "Consultants + corporate", RSM_GRAY)
    add_metric(slide, 8.05, 2.25, 2.2, 1.25, "Student funnel", "50k", "Research adoption layer", RSM_BLUE)
    add_metric(slide, 10.5, 2.25, 2.2, 1.25, "FY5 penetration", "1.5%", "Workflow seats / corporate seats", RSM_GREEN)
    make_table_ppt(
        slide,
        0.7,
        4.05,
        11.95,
        1.65,
        ["Segment", "FY5 unit target", "Penetration implied", "What must be proven"],
        [
            ["Lite", "8,000 seats", "16.0% of students/academics", "Low-cost acquisition without support overload"],
            ["Premium", "2,000 seats", "4.9% of professional universe", "Paid research value vs DDTC/Hukumonline/Ortax"],
            ["Workflow", "500 seats", "1.5% of corporate seats", "Dispute drafting productivity and accuracy"],
            ["Enterprise", "30 accounts", "0.135% of companies with tax function", "Security, integration, and firm-level governance"],
        ],
        [1.2, 1.4, 1.5, 4.2],
        font_size=7,
    )
    add_small_footer(slide)

    # 4 Solution architecture
    slide = prs.slides.add_slide(blank)
    add_title(slide, "3. Product concept: agentic tax-dispute operating layer", "A stronger proposition is a workflow that turns documents + regulations + precedents into reviewable workpapers.")
    stages = [
        ("Ingest", "PDF decisions, regulations, client documents"),
        ("Extract", "Parties, tax type, issue, values, evidence, outcome"),
        ("Retrieve", "Comparable decisions and relevant rules"),
        ("Analyze", "Risk, evidence gap, argument fit, scorecard"),
        ("Draft", "Word/PDF report for advisor review"),
    ]
    x = 0.75
    for i, (t, b) in enumerate(stages):
        add_card(slide, x + i * 2.48, 2.35, 2.08, 2.15, t, b, [RSM_BLUE, RSM_GREEN, RSM_GRAY, RSM_BLUE, RSM_GREEN][i], title_size=15)
        if i < len(stages) - 1:
            add_text(slide, "→", x + i * 2.48 + 2.15, 3.18, 0.3, 0.3, size=20, bold=True, color=RSM_MUTED)
    add_card(slide, 1.3, 5.25, 10.7, 0.75, "Management implication", "The product must be sold as a controlled professional workflow with source traceability, not as an autonomous legal/tax opinion engine.", RSM_GREEN, title_size=13)
    add_small_footer(slide)

    # 5 Business model
    slide = prs.slides.add_slide(blank)
    add_title(slide, "4. Business model: hybrid SaaS + workflow + enterprise", "Low-price research can create the funnel, but the economic engine is workflow and enterprise.")
    make_table_ppt(
        slide,
        0.75,
        2.2,
        6.0,
        3.5,
        ["Tier", "Price / year", "Role in funnel"],
        [[t, rp_mm(price), role] for t, role, price, _ in PRICING],
        [1.1, 1.2, 3.2],
        font_size=8,
    )
    make_table_ppt(
        slide,
        7.15,
        2.2,
        5.35,
        3.5,
        ["Tier", "LTV/CAC", "Payback"],
        [[u[0], u[6], u[7]] for u in UNIT_ECON],
        [1.3, 1.0, 1.4],
        font_size=8,
    )
    add_card(slide, 0.75, 6.0, 11.75, 0.65, "Key recommendation", "Do not over-index on Lite revenue. Use Lite for education and lead generation; protect margin through Premium, Workflow, and Enterprise conversion.", RSM_BLUE, title_size=12)
    add_small_footer(slide)

    # 6 Financial
    slide = prs.slides.add_slide(blank)
    add_title(slide, "5. Baseline financial model: attractive, but back-loaded", "The model turns EBITDA positive in FY3; cumulative EBITDA becomes positive in FY4.")
    max_rev = max(REVENUE)
    max_eb = max(EBITDA)
    for i, y in enumerate(YEARS):
        add_bar(slide, 0.9, 2.35 + i * 0.45, y + " revenue", REVENUE[i] / 1000, max_rev / 1000, RSM_BLUE, "bn")
        add_bar(slide, 7.0, 2.35 + i * 0.45, y + " EBITDA", EBITDA[i] / 1000, max_eb / 1000, RSM_GREEN if EBITDA[i] >= 0 else "C0504D", "bn")
    add_metric(slide, 0.9, 5.35, 2.5, 1.05, "FY5 gross margin", pct(GM[-1]), "Not 50%; 50% is closer to FY4 EBITDA margin", RSM_BLUE)
    add_metric(slide, 3.7, 5.35, 2.5, 1.05, "FY5 EBITDA margin", pct(EBITDA[-1] / REVENUE[-1]), "Baseline", RSM_GREEN)
    add_metric(slide, 6.5, 5.35, 2.5, 1.05, "Cumulative EBITDA", rp_bn(sum(EBITDA)), "5-year total", RSM_GRAY)
    add_metric(slide, 9.3, 5.35, 2.5, 1.05, "Break-even revenue FY3", rp_bn(BEP_REVENUE[2]), "Vs FY3 revenue Rp6.84bn", RSM_BLUE)
    add_small_footer(slide, "Correction: the original deck should distinguish annual EBITDA profitability, cumulative EBITDA payback, and gross margin.")

    # 7 Scenarios
    slide = prs.slides.add_slide(blank)
    add_title(slide, "6. Scenario analysis: transparent replacement for broken #REF! formulas", "Revenue scales by volume multiplier; COGS scales with gross margin; opex is held sticky.")
    make_table_ppt(
        slide,
        0.7,
        2.2,
        12.0,
        3.1,
        ["Scenario", "Multiplier", "FY5 revenue", "FY5 EBITDA", "NPV @15%", "IRR", "Payback"],
        [
            [
                r["name"],
                f"{r['multiplier']:.1f}x",
                rp_bn(r["fy5_revenue"]),
                rp_bn(r["fy5_ebitda"]),
                rp_bn(r["npv15"]),
                pct(r["irr"], 1),
                f"{r['payback']:.2f} yrs" if r["payback"] else "n/a",
            ]
            for r in scenario_rows()
        ],
        [1.35, 1.0, 1.2, 1.2, 1.15, 0.9, 1.0],
        font_size=9,
    )
    add_card(slide, 1.0, 5.75, 5.35, 0.8, "Downside insight", "Worst-case becomes value-destructive at NPV -Rp2.5bn and IRR 4.3%; paid conversion is the key sensitivity.", RSM_GRAY, title_size=12)
    add_card(slide, 6.85, 5.75, 5.35, 0.8, "Upside insight", "Moderate case reaches NPV Rp16.1bn and IRR 56.9%; valuation is highly sensitive to FY4-FY5 scale.", RSM_GREEN, title_size=12)
    add_small_footer(slide)

    # 8 Differentiation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "7. Competitive advantage: avoid the commodity research trap", "The moat should be proprietary dispute data + reviewer workflow + source traceability.")
    add_card(slide, 0.75, 2.25, 3.65, 3.4, "Commodity layer", "Tax regulation search and generic AI Q&A are easy to copy and likely face pricing pressure.", RSM_GRAY)
    add_card(slide, 4.85, 2.25, 3.65, 3.4, "Defensible layer", "Labeled outcomes, comparable decision analysis, evidence gap scoring, and advisor-reviewed templates.", RSM_BLUE)
    add_card(slide, 8.95, 2.25, 3.65, 3.4, "RSM advantage", "Brand trust, specialist reviewers, internal case experience, and enterprise client relationships.", RSM_GREEN)
    add_small_footer(slide)

    # 9 Risks
    slide = prs.slides.add_slide(blank)
    add_title(slide, "8. Risk and governance: make control part of the product", "For tax disputes, user trust depends on traceability and human review.")
    make_table_ppt(
        slide,
        0.75,
        2.15,
        12.0,
        3.7,
        ["Risk", "Likelihood x impact", "Mitigation in product / operating model"],
        [
            ["Privacy and taxpayer confidentiality", "High x High", "Data wall, anonymization, RBAC, audit log, PDP/legal review."],
            ["Wrong legal/tax inference", "Medium x High", "RAG citations, confidence scoring, reviewer approval, source-open links."],
            ["Data quality and extraction error", "High x Medium", "Validation workflow, re-extraction, field completeness score, reviewer notes."],
            ["Weak paid conversion", "Medium x Medium", "Pilot KPIs, sales qualification, premium feature gating."],
            ["Capital/model credibility", "Medium x High", "Assumption register, sensitivity table, board-ready model controls."],
        ],
        [2.2, 1.4, 4.9],
        font_size=8,
    )
    add_small_footer(slide, "Regulatory context: Indonesia UU No. 27/2022 on personal data protection should be part of the governance checklist.")

    # 10 Decision ask
    slide = prs.slides.add_slide(blank)
    add_title(slide, "9. Recommended decision path", "Approve staged validation before committing to full build-out.")
    add_card(slide, 0.75, 2.2, 3.7, 3.45, "0-90 days", "1) Repair financial model\n2) Internal RSM pilot\n3) Validate source data\n4) Build KPI dashboard\n5) 20-30 WTP interviews", RSM_BLUE)
    add_card(slide, 4.85, 2.2, 3.7, 3.45, "90-180 days", "1) Paid beta Premium/Workflow\n2) Security and governance review\n3) Expand decision labels\n4) Pricing experiments\n5) Partner-channel test", RSM_GREEN)
    add_card(slide, 8.95, 2.2, 3.7, 3.45, "Decision gates", "Proceed only if pilot proves time saving, draft quality, paid willingness, and source traceability.", RSM_GRAY)
    add_small_footer(slide)

    # 11 Appendix slide fixes
    slide = prs.slides.add_slide(blank)
    add_title(slide, "10. Slide refinements inserted in this version", "What changed versus the source deck.")
    add_bullets_ppt(
        slide,
        [
            "Reframed the proposition from 'AI tax library' to 'agentic tax-dispute workflow'.",
            "Added corrected financial scorecard: FY5 revenue, FY5 EBITDA, NPV, IRR, payback.",
            "Replaced broken scenario sheet with transparent worst/basic/baseline/moderate cases.",
            "Corrected margin language: FY5 gross margin is 80%; FY4 EBITDA margin is close to 50%.",
            "Added explicit governance risk slide covering privacy, hallucination, source traceability, and reviewer controls.",
            "Added decision-gate roadmap so leadership can approve a staged pilot instead of a one-shot large commitment.",
        ],
        0.95,
        2.2,
        11.4,
        3.8,
        size=14,
    )
    add_small_footer(slide)

    prs.save(PPTX_OUT)


if __name__ == "__main__":
    build_docx()
    build_pptx()
    print(f"Wrote {DOCX_OUT}")
    print(f"Wrote {PPTX_OUT}")
